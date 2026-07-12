"""Local study corpus: download + text-extract + index + full-text search.

Layout (under config.CACHE_DIR, default ~/.moodle-mcp):
    <courseid>/
        index.json          # {courseid, fullname, items: {cid: record}}
        files/<filename>    # downloaded originals
        text/<cid>.txt      # extracted plain text
    search.db               # SQLite FTS5 index across ALL courses

Incremental sync (module-type aware):

    core_course_get_contents
            |
      section -> module ---- modname/type branch:
            |                    file/resource  -> download if timemodified newer, extract
            |                    url            -> store external link, DO NOT download
            |                    page/label     -> extract text from description/index.html
            v
    write index.json, prune items no longer upstream, rebuild FTS rows for course

Change detection is by ``timemodified`` (free metadata from the API), never by
hashing — hashing would require downloading the file first, defeating the point.
"""
from __future__ import annotations

import html
import json
import re
import sqlite3
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

from . import classify, config
from .client import MoodleAPIError, MoodleClient

_UNSAFE = re.compile(r'[<>:"/\\|?*\x00-\x1f]')


def _safe_filename(name: str) -> str:
    """Strip path separators / unsafe chars so a Moodle filename can't escape the dir."""
    name = name.replace("\\", "/").split("/")[-1]
    name = _UNSAFE.sub("_", name).strip(". ")
    return name or "unnamed"


def _file_dest(files_dir: Path, cid: str, filename: str) -> Path:
    """On-disk path for a file, namespaced by module id so identically-named
    files in different modules never collide (or race under concurrency)."""
    cmid = cid.split(":", 1)[0]
    return files_dir / f"{cmid}__{_safe_filename(filename)}"


@dataclass
class Item:
    cid: str
    name: str
    category: str
    confidence: str
    kind: str  # "file" | "url" | "page"
    modname: str = ""
    section: str = ""
    timemodified: int = 0
    filesize: int = 0
    filename: str = ""
    fileurl: str = ""
    external_url: str = ""
    text_extracted: bool = False
    text_chars: int = 0


@dataclass
class SyncResult:
    courseid: int
    fullname: str
    downloaded: int = 0
    skipped: int = 0
    pruned: int = 0
    links: int = 0
    extract_failures: list[str] = field(default_factory=list)
    errors: list[str] = field(default_factory=list)
    items: int = 0


# -- text extraction -------------------------------------------------------
def _strip_html(html: str) -> str:
    text = re.sub(r"(?is)<(script|style)[^>]*>.*?</\1>", " ", html)
    text = re.sub(r"(?s)<[^>]+>", " ", text)
    text = re.sub(r"&nbsp;", " ", text)
    text = re.sub(r"&amp;", "&", text)
    text = re.sub(r"&lt;", "<", text)
    text = re.sub(r"&gt;", ">", text)
    return re.sub(r"\s+", " ", text).strip()


def extract_text(path: Path) -> tuple[str, bool]:
    """Return (text, ok). ok=False for scanned PDFs / unsupported / failures."""
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            parts = [(page.extract_text() or "") for page in reader.pages]
            text = "\n".join(parts).strip()
            return text, bool(text)
        if suffix in (".txt", ".md", ".csv", ".tex"):
            text = path.read_text(encoding="utf-8", errors="replace").strip()
            return text, bool(text)
        if suffix in (".html", ".htm"):
            text = _strip_html(path.read_text(encoding="utf-8", errors="replace"))
            return text, bool(text)
        if suffix == ".docx":
            try:
                import docx  # type: ignore

                text = "\n".join(p.text for p in docx.Document(str(path)).paragraphs).strip()
                return text, bool(text)
            except ImportError:
                return "", False
        if suffix == ".pptx":
            try:
                from pptx import Presentation  # type: ignore

                chunks: list[str] = []
                for slide in Presentation(str(path)).slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            chunks.append(shape.text_frame.text)
                text = "\n".join(chunks).strip()
                return text, bool(text)
            except ImportError:
                return "", False
    except Exception:  # noqa: BLE001 - extraction must never crash a sync
        return "", False
    return "", False


# -- paths & index ---------------------------------------------------------
def _course_dir(courseid: int) -> Path:
    return config.CACHE_DIR / str(courseid)


def _index_path(courseid: int) -> Path:
    return _course_dir(courseid) / "index.json"


def load_index(courseid: int) -> dict[str, Any]:
    path = _index_path(courseid)
    if path.exists():
        try:
            return json.loads(path.read_text(encoding="utf-8"))
        except (ValueError, OSError):
            pass
    return {"courseid": courseid, "fullname": "", "items": {}}


def _save_index(courseid: int, fullname: str, items: dict[str, Item]) -> None:
    path = _index_path(courseid)
    path.parent.mkdir(parents=True, exist_ok=True)
    payload = {
        "courseid": courseid,
        "fullname": fullname,
        "items": {cid: asdict(it) for cid, it in items.items()},
    }
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")


# -- FTS index -------------------------------------------------------------
def _connect_db() -> tuple[sqlite3.Connection, bool]:
    """Open the search DB. Returns (conn, fts5) — fts5=False means LIKE fallback."""
    config.CACHE_DIR.mkdir(parents=True, exist_ok=True)
    conn = sqlite3.connect(config.CACHE_DIR / "search.db")
    try:
        conn.execute(
            "CREATE VIRTUAL TABLE IF NOT EXISTS materials USING fts5("
            "courseid UNINDEXED, cid UNINDEXED, name, category UNINDEXED, body)"
        )
        return conn, True
    except sqlite3.OperationalError:
        conn.execute(
            "CREATE TABLE IF NOT EXISTS materials ("
            "courseid INTEGER, cid TEXT, name TEXT, category TEXT, body TEXT)"
        )
        return conn, False


def _reindex_course(courseid: int, rows: list[tuple[int, str, str, str, str]]) -> None:
    conn, _ = _connect_db()
    try:
        conn.execute("DELETE FROM materials WHERE courseid = ?", (courseid,))
        conn.executemany(
            "INSERT INTO materials (courseid, cid, name, category, body) VALUES (?,?,?,?,?)",
            rows,
        )
        conn.commit()
    finally:
        conn.close()


def search(query: str, courseid: int | None = None, limit: int = 20) -> list[dict[str, Any]]:
    """Full-text search across synced materials. courseid=None searches all."""
    conn, fts5 = _connect_db()
    try:
        if fts5:
            sql = (
                "SELECT courseid, cid, name, category, snippet(materials, 4, '[', ']', ' … ', 12) "
                "FROM materials WHERE materials MATCH ?"
            )
            params: list[Any] = [query]
            if courseid is not None:
                sql += " AND courseid = ?"
                params.append(courseid)
            sql += " LIMIT ?"
            params.append(limit)
            try:
                cur = conn.execute(sql, params)
            except sqlite3.OperationalError:
                # malformed MATCH expr -> fall back to a quoted phrase search
                params[0] = '"' + query.replace('"', " ") + '"'
                cur = conn.execute(sql, params)
        else:
            sql = (
                "SELECT courseid, cid, name, category, substr(body,1,240) "
                "FROM materials WHERE body LIKE ?"
            )
            params = [f"%{query}%"]
            if courseid is not None:
                sql += " AND courseid = ?"
                params.append(courseid)
            sql += " LIMIT ?"
            params.append(limit)
            cur = conn.execute(sql, params)
        return [
            {"courseid": r[0], "cid": r[1], "name": r[2], "category": r[3], "snippet": r[4]}
            for r in cur.fetchall()
        ]
    finally:
        conn.close()


def get_material_text(courseid: int, cid: str) -> str | None:
    path = _course_dir(courseid) / "text" / f"{_safe_filename(cid)}.txt"
    if path.exists():
        return path.read_text(encoding="utf-8", errors="replace")
    return None


# -- sync ------------------------------------------------------------------
def _write_text(text_dir: Path, cid: str, text: str) -> None:
    text_dir.mkdir(parents=True, exist_ok=True)
    (text_dir / f"{_safe_filename(cid)}.txt").write_text(text, encoding="utf-8")


def _download_worker(client: MoodleClient, task: dict[str, Any], text_dir: Path):
    """Thread worker: download one file + extract text. Returns (Item|None, status, info)."""
    try:
        client.download_file(task["fileurl"], task["dest"])
    except MoodleAPIError as exc:
        return None, "error", f"{task['filename']}: {exc}"
    text, ok = extract_text(task["dest"])
    if ok:
        _write_text(text_dir, task["cid"], text)
    item = Item(
        cid=task["cid"], name=task["display"], category=task["category"],
        confidence=task["confidence"], kind="file", modname=task["modname"],
        section=task["section"], timemodified=task["timemodified"],
        filesize=task["filesize"], filename=task["filename"], fileurl=task["fileurl"],
        text_extracted=ok, text_chars=len(text),
    )
    return item, ("ok" if ok else "extract_fail"), task["filename"]


def sync_course(
    client: MoodleClient, courseid: int, fullname: str = "", *, max_workers: int = 6
) -> SyncResult:
    """Download + extract + index a course. Incremental by timemodified.

    Two phases: walk the contents (reusing unchanged files, queuing changed ones),
    then download the queue concurrently — a first sync of a large course is IO
    bound, so parallel GETs turn minutes into tens of seconds.
    """
    result = SyncResult(courseid=courseid, fullname=fullname)
    contents = client.get_course_contents(courseid)

    prev_items = load_index(courseid).get("items", {})
    items: dict[str, Item] = {}
    course_dir = _course_dir(courseid)
    files_dir = course_dir / "files"
    text_dir = course_dir / "text"

    pending: list[dict[str, Any]] = []
    for section in contents:
        section_name = html.unescape(section.get("name", ""))
        for module in section.get("modules", []):
            modname = module.get("modname", "")
            name = html.unescape(module.get("name", ""))
            cmid = module.get("id", 0)
            cls = classify.classify_item(name, "", section_name, modname)
            contents_list = module.get("contents", []) or []
            file_contents = [c for c in contents_list if c.get("type") == "file"]
            url_contents = [c for c in contents_list if c.get("type") == "url"]

            if file_contents:
                multi = len(file_contents) > 1
                for c in file_contents:
                    filename = c.get("filename", "file")
                    cid = f"{cmid}:{filename}"
                    timemodified = int(c.get("timemodified", 0) or 0)
                    filesize = int(c.get("filesize", 0) or 0)
                    fileurl = c.get("fileurl", "")
                    dest = _file_dest(files_dir, cid, filename)
                    prev = prev_items.get(cid)
                    # Classify per file (by filename) so a folder of exams counts
                    # each exam, not the folder as one item.
                    fcls = classify.classify_item(name, filename, section_name, modname)
                    display = filename if multi else (name or filename)

                    if (
                        prev
                        and int(prev.get("timemodified", -1)) == timemodified
                        and dest.exists()
                    ):
                        items[cid] = Item(**{k: prev[k] for k in Item.__dataclass_fields__ if k in prev})
                        result.skipped += 1
                        continue

                    pending.append({
                        "cid": cid, "fileurl": fileurl, "dest": dest, "display": display,
                        "category": fcls.category, "confidence": fcls.confidence,
                        "modname": modname, "section": section_name,
                        "timemodified": timemodified, "filesize": filesize, "filename": filename,
                    })

            elif url_contents:
                cid = str(cmid)
                items[cid] = Item(
                    cid=cid, name=name, category=cls.category, confidence=cls.confidence,
                    kind="url", modname=modname, section=section_name,
                    external_url=url_contents[0].get("fileurl", ""),
                )
                result.links += 1

            elif module.get("description"):
                cid = str(cmid)
                text = _strip_html(module["description"])
                if text:
                    _write_text(text_dir, cid, text)
                items[cid] = Item(
                    cid=cid, name=name, category=cls.category, confidence=cls.confidence,
                    kind="page", modname=modname, section=section_name,
                    text_extracted=bool(text), text_chars=len(text),
                )

    # Phase 2: download changed/new files concurrently (IO-bound).
    if pending:
        files_dir.mkdir(parents=True, exist_ok=True)
        with ThreadPoolExecutor(max_workers=max_workers) as pool:
            futures = [pool.submit(_download_worker, client, t, text_dir) for t in pending]
            for fut in as_completed(futures):
                item, status, info = fut.result()
                if item is None:
                    result.errors.append(info)
                    continue
                items[item.cid] = item
                result.downloaded += 1
                if status == "extract_fail":
                    result.extract_failures.append(info)

    # Prune orphaned items: delete their cached file + extracted text.
    removed = set(prev_items) - set(items)
    for cid in removed:
        prev = prev_items.get(cid, {})
        fn = prev.get("filename")
        if fn:
            fp = _file_dest(files_dir, cid, fn)
            if fp.exists():
                try:
                    fp.unlink()
                except OSError:
                    pass
        tp = text_dir / f"{_safe_filename(cid)}.txt"
        if tp.exists():
            try:
                tp.unlink()
            except OSError:
                pass
    result.pruned = len(removed)

    result.items = len(items)
    _save_index(courseid, fullname, items)

    rows: list[tuple[int, str, str, str, str]] = [
        (courseid, cid, it.name, it.category, get_material_text(courseid, cid) or "")
        for cid, it in items.items()
    ]
    _reindex_course(courseid, rows)

    return result
